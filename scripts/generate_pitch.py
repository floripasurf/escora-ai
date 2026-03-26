"""Generate Escora.AI pitch presentation (PPTX)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# Brand colors
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x16, 0x84, 0x7D)  # teal
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)

WIDTH = Inches(13.333)
HEIGHT = Inches(7.5)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text(slide, left, top, width, height, text, font_size=18,
              color=DARK, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def _add_bullet_list(slide, left, top, width, height, items, font_size=16,
                     color=DARK, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def _add_metric_box(slide, left, top, width, height, value, label,
                    value_color=ACCENT):
    shape = _add_shape(slide, left, top, width, height, WHITE)
    shape.shadow.inherit = False

    # Value
    _add_text(slide, left, top + Inches(0.2), width, Inches(0.6),
              value, font_size=36, color=value_color, bold=True,
              align=PP_ALIGN.CENTER)
    # Label
    _add_text(slide, left, top + Inches(0.85), width, Inches(0.5),
              label, font_size=12, color=MEDIUM_GRAY,
              align=PP_ALIGN.CENTER)


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, DARK)

    # Accent bar top
    _add_shape(slide, Inches(0), Inches(0), WIDTH, Inches(0.08), ACCENT)

    # Title
    _add_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
              "Escora.AI", font_size=60, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER)

    # Subtitle
    _add_text(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
              "Cálculo automatizado de escoramento a partir de projetos DXF",
              font_size=24, color=RGBColor(0xAA, 0xAA, 0xAA),
              align=PP_ALIGN.CENTER)

    # Bottom line
    _add_shape(slide, Inches(5.5), Inches(4.2), Inches(2.3), Inches(0.04), ACCENT)

    # Footer
    _add_text(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
              "Florianópolis, 2026", font_size=14,
              color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), Inches(0.08), HEIGHT, ACCENT)

    _add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
              "O Problema", font_size=32, color=DARK, bold=True)
    _add_shape(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT)

    problems = [
        "O cálculo de escoramento para formas de concreto é feito\n"
        "manualmente por engenheiros — processo lento, repetitivo e\n"
        "sujeito a erros",
        "Cada projeto de fôrma exige: identificação de vigas e pilares,\n"
        "cálculo de cargas, seleção de escoras, definição de espaçamentos\n"
        "e geração de relatórios",
        "Um engenheiro gasta 4–8 horas por pavimento;\n"
        "um prédio de 20 andares pode levar semanas",
        "Erros de cálculo podem causar acidentes graves na obra\n"
        "(NBR 15696 — Segurança de estruturas provisórias)",
    ]

    y = Inches(1.5)
    for i, text in enumerate(problems):
        # Number circle
        num_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.05),
            Inches(0.4), Inches(0.4))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = ACCENT if i < 3 else ORANGE
        num_shape.line.fill.background()
        tf = num_shape.text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        _add_text(slide, Inches(1.5), y, Inches(10), Inches(1.0),
                  text, font_size=16, color=DARK)
        y += Inches(1.3)


def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), Inches(0.08), HEIGHT, ACCENT)

    _add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
              "A Solução", font_size=32, color=DARK, bold=True)
    _add_shape(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT)

    _add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
              "O engenheiro envia o arquivo DXF do projeto estrutural.\n"
              "O Escora.AI faz o resto automaticamente:",
              font_size=18, color=MEDIUM_GRAY)

    steps = [
        ("Leitura do DXF", "Interpreta geometria, layers e anotações\n"
         "do projeto estrutural (vigas, pilares, lajes)"),
        ("Classificação", "Identifica elementos por geometria (pares paralelos,\n"
         "retângulos) e texto (V1, P3, 14x40) com score de confiança"),
        ("Cálculo", "Calcula cargas (NBR 6120), seleciona escoras do\n"
         "catálogo, distribui espaçamentos, valida capacidades"),
        ("Relatório", "Gera PDF profissional + planilha Excel com\n"
         "BOM (lista de materiais), pronto para compra"),
    ]

    x_start = Inches(0.8)
    y = Inches(2.5)
    box_w = Inches(2.8)
    gap = Inches(0.3)

    for i, (title, desc) in enumerate(steps):
        x = x_start + i * (box_w + gap)

        # Box background
        _add_shape(slide, x, y, box_w, Inches(3.2), LIGHT_GRAY)

        # Step number
        _add_text(slide, x, y + Inches(0.15), box_w, Inches(0.4),
                  f"0{i+1}", font_size=28, color=ACCENT, bold=True,
                  align=PP_ALIGN.CENTER)

        # Title
        _add_text(slide, x + Inches(0.2), y + Inches(0.7), box_w - Inches(0.4), Inches(0.4),
                  title, font_size=18, color=DARK, bold=True,
                  align=PP_ALIGN.CENTER)

        # Description
        _add_text(slide, x + Inches(0.2), y + Inches(1.2), box_w - Inches(0.4), Inches(1.8),
                  desc, font_size=13, color=MEDIUM_GRAY,
                  align=PP_ALIGN.CENTER)

        # Arrow between boxes
        if i < 3:
            arrow_x = x + box_w + Inches(0.05)
            _add_text(slide, arrow_x, y + Inches(1.2), Inches(0.2), Inches(0.4),
                      "→", font_size=24, color=ACCENT, bold=True,
                      align=PP_ALIGN.CENTER)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), Inches(0.08), HEIGHT, ACCENT)

    _add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
              "Arquitetura do Pipeline", font_size=32, color=DARK, bold=True)
    _add_shape(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT)

    stages = [
        ("Stage 1\nParse", "Leitura do DXF\n(ezdxf)", ACCENT),
        ("Stage 2\nSegment", "Separação por\nnível/pavimento", ACCENT),
        ("Stage 3\nClassify", "Geometria +\nTexto → Score", ACCENT),
        ("Stage 4\nMetadata", "Pé-direito\nEspessura", ACCENT),
        ("Stage 5\nCalculate", "Cargas, escoras\nespaçamentos", ACCENT),
        ("Stage 6\nOutput", "PDF + Excel\nRelatório", ORANGE),
    ]

    y = Inches(2.0)
    box_w = Inches(1.7)
    box_h = Inches(2.0)
    gap = Inches(0.25)
    x_start = Inches(0.8)

    for i, (title, desc, color) in enumerate(stages):
        x = x_start + i * (box_w + gap)

        shape = _add_shape(slide, x, y, box_w, box_h, color)

        _add_text(slide, x, y + Inches(0.2), box_w, Inches(0.7),
                  title, font_size=14, color=WHITE, bold=True,
                  align=PP_ALIGN.CENTER)

        _add_text(slide, x, y + Inches(0.9), box_w, Inches(0.9),
                  desc, font_size=12, color=WHITE,
                  align=PP_ALIGN.CENTER)

        if i < 5:
            arrow_x = x + box_w + Inches(0.02)
            _add_text(slide, arrow_x, y + Inches(0.7), Inches(0.2), Inches(0.4),
                      "→", font_size=20, color=ACCENT, bold=True)

    # Tech stack
    _add_text(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.4),
              "Stack Tecnológico", font_size=18, color=DARK, bold=True)

    techs = [
        ("Python 3.11+", "Linguagem principal"),
        ("ezdxf", "Leitura de arquivos DXF"),
        ("Shapely", "Geometria computacional"),
        ("ReportLab", "Geração de PDF"),
        ("openpyxl", "Geração de Excel"),
        ("FastAPI", "API REST"),
    ]

    x = Inches(0.8)
    for i, (tech, desc) in enumerate(techs):
        col = i % 3
        row = i // 3
        tx = x + col * Inches(3.8)
        ty = Inches(5.0) + row * Inches(0.7)

        _add_text(slide, tx, ty, Inches(1.5), Inches(0.35),
                  tech, font_size=14, color=ACCENT, bold=True)
        _add_text(slide, tx + Inches(1.6), ty, Inches(2.0), Inches(0.35),
                  desc, font_size=13, color=MEDIUM_GRAY)


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK)
    _add_shape(slide, Inches(0), Inches(0), WIDTH, Inches(0.08), ACCENT)

    _add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
              "Resultados Reais", font_size=32, color=WHITE, bold=True)
    _add_shape(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT)

    _add_text(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.4),
              "Projeto CFL-SUB — Subsolo com 83+ pilares e 70+ vigas",
              font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA))

    # Metric boxes
    metrics = [
        ("71", "Vigas\ndetectadas"),
        ("99", "Pilares\ndetectados"),
        ("618", "Escoras\ncalculadas"),
        ("6.190 kN", "Carga\ntotal"),
        ("< 10s", "Tempo de\nprocessamento"),
    ]

    y = Inches(2.2)
    box_w = Inches(2.1)
    box_h = Inches(1.5)
    gap = Inches(0.3)
    x_start = Inches(0.8)

    for i, (value, label) in enumerate(metrics):
        x = x_start + i * (box_w + gap)
        shape = _add_shape(slide, x, y, box_w, box_h,
                           RGBColor(0x22, 0x22, 0x3A))

        _add_text(slide, x, y + Inches(0.2), box_w, Inches(0.5),
                  value, font_size=32, color=ACCENT, bold=True,
                  align=PP_ALIGN.CENTER)
        _add_text(slide, x, y + Inches(0.8), box_w, Inches(0.5),
                  label, font_size=12, color=RGBColor(0xAA, 0xAA, 0xAA),
                  align=PP_ALIGN.CENTER)

    # Deliverables
    _add_text(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.4),
              "Entregáveis gerados automaticamente:", font_size=18,
              color=WHITE, bold=True)

    deliverables = [
        "PDF — Relatório de escoramento com tabelas de vigas, lajes e BOM",
        "Excel — Planilha com 3 abas (BOM, Vigas, Lajes) para orçamento",
        "Validação — Alertas de escoras sobrecarregadas e valores padrão usados",
    ]

    _add_bullet_list(slide, Inches(1.0), Inches(4.8), Inches(10), Inches(2.0),
                     [f"•  {d}" for d in deliverables], font_size=15,
                     color=RGBColor(0xCC, 0xCC, 0xCC))


def slide_compliance(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), Inches(0.08), HEIGHT, ACCENT)

    _add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
              "Conformidade Normativa", font_size=32, color=DARK, bold=True)
    _add_shape(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT)

    norms = [
        ("NBR 15696:2009", "Fôrmas e escoramentos para estruturas de concreto\n"
         "— Projeto, dimensionamento e procedimentos executivos"),
        ("NBR 6120:2019", "Ações para o cálculo de estruturas de edificações\n"
         "— Cargas permanentes e acidentais"),
        ("NBR 6118:2023", "Projeto de estruturas de concreto\n"
         "— Procedimento (referência para seções de vigas)"),
    ]

    y = Inches(1.6)
    for code, desc in norms:
        _add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.2), LIGHT_GRAY)
        _add_text(slide, Inches(1.0), y + Inches(0.15), Inches(3), Inches(0.4),
                  code, font_size=20, color=ACCENT, bold=True)
        _add_text(slide, Inches(1.0), y + Inches(0.55), Inches(10), Inches(0.6),
                  desc, font_size=14, color=MEDIUM_GRAY)
        y += Inches(1.5)

    # Safety note
    _add_text(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.8),
              "O sistema valida automaticamente se as escoras selecionadas\n"
              "atendem à capacidade de carga exigida e alerta quando há sobrecarga.",
              font_size=16, color=DARK)


def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), Inches(0.08), HEIGHT, ACCENT)

    _add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
              "Roadmap", font_size=32, color=DARK, bold=True)
    _add_shape(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT)

    phases = [
        ("Atual ✓", [
            "Pipeline DXF → Cálculo → PDF/Excel",
            "Detecção de vigas, pilares e lajes",
            "Seleção automática de escoras do catálogo",
            "Relatório profissional em português",
        ], ACCENT),
        ("Próximo", [
            "DXF output — posições de escoras no projeto",
            "Catálogo real de fabricantes (SH, Metax)",
            "Refinamento de detecção com biblioteca DXF",
            "Interface web para upload e visualização",
        ], RGBColor(0x33, 0x77, 0xBB)),
        ("Futuro", [
            "Branding por cliente (logo no relatório)",
            "Múltiplos pavimentos simultâneos",
            "Integração com ERPs de construção",
            "Marketplace de catálogos de escoras",
        ], MEDIUM_GRAY),
    ]

    x_start = Inches(0.8)
    box_w = Inches(3.6)
    gap = Inches(0.4)
    y = Inches(1.6)

    for i, (title, items, color) in enumerate(phases):
        x = x_start + i * (box_w + gap)

        # Header
        _add_shape(slide, x, y, box_w, Inches(0.5), color)
        _add_text(slide, x, y + Inches(0.05), box_w, Inches(0.4),
                  title, font_size=16, color=WHITE, bold=True,
                  align=PP_ALIGN.CENTER)

        # Items
        _add_shape(slide, x, y + Inches(0.5), box_w, Inches(4.0), LIGHT_GRAY)
        _add_bullet_list(slide, x + Inches(0.2), y + Inches(0.7),
                         box_w - Inches(0.4), Inches(3.5),
                         [f"•  {item}" for item in items],
                         font_size=13, color=DARK, spacing=Pt(12))


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK)
    _add_shape(slide, Inches(0), Inches(0), WIDTH, Inches(0.08), ACCENT)

    _add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
              "Escora.AI", font_size=54, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER)

    _add_text(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.6),
              "Cálculo de escoramento inteligente.\nDo DXF ao relatório em segundos.",
              font_size=22, color=RGBColor(0xAA, 0xAA, 0xAA),
              align=PP_ALIGN.CENTER)

    _add_shape(slide, Inches(5.5), Inches(4.2), Inches(2.3), Inches(0.04), ACCENT)

    _add_text(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
              "contato@escora.ai", font_size=16,
              color=ACCENT, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = WIDTH
    prs.slide_height = HEIGHT

    slide_cover(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_architecture(prs)
    slide_results(prs)
    slide_compliance(prs)
    slide_roadmap(prs)
    slide_closing(prs)

    output = Path("output/Escora.AI-Apresentacao.pptx")
    output.parent.mkdir(exist_ok=True)
    prs.save(str(output))
    print(f"Apresentação salva: {output}")
    print(f"Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
