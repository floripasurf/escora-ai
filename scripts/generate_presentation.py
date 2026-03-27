"""Generate Escora.AI presentation in PPTX format."""
import sys
sys.path.insert(0, "/Users/raphaellages/Desktop/escora-ai")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)       # Dark navy
BG_MEDIUM = RGBColor(0x16, 0x21, 0x3E)      # Medium navy
ACCENT_ORANGE = RGBColor(0xE8, 0x7D, 0x2F)  # Construction orange
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)    # Tech blue
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)   # Success green
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)     # Alert red
ACCENT_YELLOW = RGBColor(0xFB, 0xBF, 0x24)  # Warning yellow
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)
DARK_CARD = RGBColor(0x1E, 0x29, 0x3B)


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def set_text(shape, text, size=18, color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_paragraph(tf, text, size=16, color=TEXT_LIGHT, bold=False, space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(space_before)
    return p


def card(slide, left, top, width, height, title, items, title_color=ACCENT_ORANGE, item_size=14):
    shape = add_shape(slide, left, top, width, height, DARK_CARD, radius=0.05)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(16)
    tf.margin_right = Pt(16)
    tf.margin_top = Pt(12)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.color.rgb = title_color
    p.font.bold = True
    for item in items:
        add_paragraph(tf, item, size=item_size, color=TEXT_LIGHT, space_before=4)
    return shape


def metric_card(slide, left, top, width, height, value, label, color=ACCENT_BLUE):
    shape = add_shape(slide, left, top, width, height, DARK_CARD, radius=0.05)
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(4), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(20)
    tf.margin_top = Pt(10)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(32)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_MUTED
    p2.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Orange accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_ORANGE
bar.line.fill.background()

# Title
tx = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(1.5))
tf = set_text(tx, "ESCORA.AI", size=60, color=TEXT_WHITE, bold=True, align=PP_ALIGN.LEFT)
add_paragraph(tf, "Inteligencia Artificial para Projetos de Escoramento", size=28, color=ACCENT_ORANGE, bold=False, space_before=8)

# Subtitle
tx2 = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(8), Inches(1.5))
tf2 = set_text(tx2, "Automacao completa: do DXF estrutural ao posicionamento otimizado de escoras", size=18, color=TEXT_LIGHT)
add_paragraph(tf2, "NBR 15696:2009 | NBR 6120:2019 | Python + Shapely + ezdxf", size=14, color=TEXT_MUTED, space_before=12)

# Version/date
tx3 = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(6), Inches(0.5))
set_text(tx3, "Marco 2026 | v0.9-beta | Raphael Lages", size=14, color=TEXT_MUTED)


# ═══════════════════════════════════════════
# SLIDE 2: The Problem
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
set_text(tx, "O PROBLEMA", size=36, color=ACCENT_ORANGE, bold=True)

card(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.2),
     "Processo Manual Atual",
     ["Engenheiro recebe planta estrutural (DXF/DWG)",
      "Identifica manualmente vigas, pilares e lajes",
      "Calcula cargas por trecho (NBR 15696)",
      "Posiciona escoras uma a uma no CAD",
      "Gera lista de materiais (BOM)",
      "Processo leva 4-8 horas por pavimento"],
     title_color=ACCENT_RED, item_size=14)

card(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.2),
     "Consequencias",
     ["Alto custo de mao-de-obra qualificada",
      "Erros humanos em calculo de cargas",
      "Retrabalho quando estoque muda",
      "Tempo de resposta lento para orcamentos",
      "Dificuldade de escalar a operacao",
      "Antiprojeto demorado = perde cliente"],
     title_color=ACCENT_RED, item_size=14)

card(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.5),
     "Oportunidade de Mercado",
     ["Brasil: >200 locadoras de escoramento ativas",
      "Cada locadora processa 5-20 projetos/semana",
      "Taxa de consumo tipica: ~10 kg/m2 (antiprojeto competitivo)",
      "Fluxo: Prospeccao > Antiprojeto > Proposta > Projeto Executivo > Assistencia Tecnica > Logistica Reversa",
      "Antiprojeto e o gargalo — precisa ser rapido e competitivo para fechar venda",
      "Engenheiros gastam 80% do tempo em tarefas repetitivas que podem ser automatizadas"],
     title_color=ACCENT_YELLOW, item_size=14)


# ═══════════════════════════════════════════
# SLIDE 3: The Solution
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
set_text(tx, "A SOLUCAO: PIPELINE DE 6 ESTAGIOS", size=36, color=ACCENT_ORANGE, bold=True)

stages = [
    ("1. PARSE", "Leitura DXF", "Extrai textos, segmentos,\nretangulos, circulos, polylines", ACCENT_BLUE),
    ("2. SEGMENT", "Segmentacao", "Agrupa entidades por\nnivel (Cobertura, Tipo...)", ACCENT_BLUE),
    ("3. CLASSIFY", "Classificacao", "Detecta vigas, pilares\nvia geometria + texto", ACCENT_ORANGE),
    ("4. METADATA", "Metadados", "Pe-direito, espessura\nda laje, escala", ACCENT_ORANGE),
    ("5. CALCULATE", "Calculo", "Cargas NBR 15696,\nposiciona escoras", ACCENT_GREEN),
    ("6. LEARN", "Aprendizado", "Salva conhecimento\npara proximas execucoes", ACCENT_GREEN),
]

for i, (num, title, desc, color) in enumerate(stages):
    left = Inches(0.5 + i * 2.1)
    top = Inches(1.8)
    shape = add_shape(slide, left, top, Inches(1.9), Inches(2.4), DARK_CARD, radius=0.05)
    # Color bar top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(1.9), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(14)
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(12)
    p.font.color.rgb = color
    p.font.bold = True
    add_paragraph(tf, title, size=16, color=TEXT_WHITE, bold=True, space_before=4)
    add_paragraph(tf, desc, size=11, color=TEXT_LIGHT, space_before=8)
    # Arrow between stages
    if i < 5:
        arrow = slide.shapes.add_textbox(left + Inches(1.95), top + Inches(1.0), Inches(0.2), Inches(0.3))
        set_text(arrow, ">", size=20, color=TEXT_MUTED, bold=True, align=PP_ALIGN.CENTER)

# Output section
card(slide, Inches(0.8), Inches(4.8), Inches(3.6), Inches(2.0),
     "Input",
     ["Arquivo DXF/DWG estrutural",
      "Qualquer escritorio de engenharia",
      "Planta de formas padrao"],
     title_color=ACCENT_BLUE, item_size=13)

card(slide, Inches(5.0), Inches(4.8), Inches(3.6), Inches(2.0),
     "Output",
     ["DXF com escoras posicionadas",
      "BOM (Lista de Materiais)",
      "Relatorio de cargas e validacao"],
     title_color=ACCENT_GREEN, item_size=13)

card(slide, Inches(9.2), Inches(4.8), Inches(3.6), Inches(2.0),
     "Aprendizado Continuo",
     ["Cada execucao acumula conhecimento",
      "Layers, secoes, pe-direito",
      "Melhora deteccao automaticamente"],
     title_color=ACCENT_ORANGE, item_size=13)


# ═══════════════════════════════════════════
# SLIDE 4: Milestones Achieved
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.8))
set_text(tx, "MILESTONES ALCANCADOS", size=36, color=ACCENT_GREEN, bold=True)

milestones = [
    ("Parser DXF Completo", "LINE, LWPOLYLINE, POLYLINE, SOLID, CIRCLE — todos os tipos de entidades estruturais"),
    ("Deteccao de Vigas", "Pares paralelos H/V com gap 8-50cm, ratio comprimento/largura >=5, deduplicacao de spans"),
    ("Deteccao de Pilares", "Retangulos (area 0.015-1.0 m2), pilares circulares (CIRCLE), cluster deduplication"),
    ("Classificacao por Layer", "Score combinado (count * rate) seleciona o melhor layer; learning boost para layers conhecidos"),
    ("Classificacao por Texto", "Regex para nomes (V1a, P5, L3), secoes (14x40, 19/60), pe-direito, espessura"),
    ("Calculo NBR 15696", "Peso proprio + sobrecarga + majoracao (gamma_f=1.4), selecao de escora por catalogo"),
    ("Distribuicao de Escoras", "Vigas: ao longo do eixo com exclusao de pilares e balanco; Lajes: grid com exclusao"),
    ("Derivacao de Lajes", "Polygonize a partir da malha de vigas; snap de endpoints para fechar gaps de pilares"),
    ("Sistema de Aprendizado", "238 records acumulados; prioriza layers historicos, secoes frequentes, pe-direito"),
    ("Output DXF", "Arquivo com layers ESCORAS_VIGA (vermelho), ESCORAS_LAJE (verde), LAJES_DET (ciano)"),
]

for i, (title, desc) in enumerate(milestones):
    col = i % 2
    row = i // 2
    left = Inches(0.6 + col * 6.2)
    top = Inches(1.3 + row * 1.15)
    shape = add_shape(slide, left, top, Inches(5.9), Inches(1.05), DARK_CARD, radius=0.03)
    # Check mark
    check = slide.shapes.add_textbox(left + Pt(8), top + Pt(8), Pt(20), Pt(20))
    set_text(check, "✓", size=14, color=ACCENT_GREEN, bold=True)
    # Text
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(30)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.font.bold = True
    add_paragraph(tf, desc, size=11, color=TEXT_MUTED, space_before=2)


# ═══════════════════════════════════════════
# SLIDE 5: Current Results (metrics)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.8))
set_text(tx, "RESULTADOS ATUAIS — CVS-COB-FOR-006-R00", size=32, color=ACCENT_ORANGE, bold=True)

# Metric cards row 1
metrics = [
    ("39", "Vigas Detectadas", ACCENT_BLUE),
    ("34", "Pilares Detectados", ACCENT_BLUE),
    ("4", "Lajes Derivadas", ACCENT_BLUE),
    ("412", "Escoras Posicionadas", ACCENT_GREEN),
    ("238", "Records de Aprendizado", ACCENT_ORANGE),
]
for i, (val, label, color) in enumerate(metrics):
    metric_card(slide, Inches(0.5 + i * 2.5), Inches(1.3), Inches(2.2), Inches(1.2), val, label, color)

# Evolution card
card(slide, Inches(0.5), Inches(2.9), Inches(5.8), Inches(3.8),
     "Evolucao da Deteccao (CVS-COB)",
     ["Rodada 1: 16 vigas, 20 pilares (spans merged, noise layers)",
      "Rodada 2: 35 vigas (BEAM_SPAN_GAP 1.0m > 0.05m)",
      "Rodada 3: 36 vigas (section normalization w < h)",
      "Rodada 4: 39 vigas (floating point fix + pillar text override)",
      "Rodada 5: 39 vigas + 4 lajes (beam grid snap para polygonize)",
      "",
      "Cada fix melhora TODOS os arquivos futuros, nao so o CVS-COB.",
      "O sistema de learning acumula conhecimento de cada rodada."],
     title_color=ACCENT_GREEN, item_size=14)

# Fixes applied card
card(slide, Inches(6.9), Inches(2.9), Inches(5.8), Inches(3.8),
     "Fixes Aplicados Nesta Iteracao",
     ["BEAM_SPAN_GAP_TOLERANCE: 1.0m > 0.05m",
      "  Vigas eram merged atraves de pilares (V15a=31m!)",
      "",
      "MIN_BEAM_LENGTH_RATIO: floating point (4.999 < 5.0)",
      "  Vigas com ratio exatamente 5.0 eram rejeitadas",
      "",
      "Pillar text override removido para vigas",
      "  Labels P20, P11 perto de vigas cancelavam deteccao",
      "",
      "Beam grid snap para derivar lajes",
      "  Endpoints das vigas estendidos ate vigas perpendiculares",
      "  Permite polygonize criar poligonos fechados de laje"],
     title_color=ACCENT_YELLOW, item_size=13)


# ═══════════════════════════════════════════
# SLIDE 6: Learning Loop (the user's idea)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.8))
set_text(tx, "LOOP DE APRENDIZADO CONTINUO", size=36, color=ACCENT_ORANGE, bold=True)

tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(0.5))
set_text(tx2, "O sistema aprende com cada revisao do engenheiro — Human-in-the-Loop AI", size=18, color=TEXT_LIGHT)

# Flow diagram as cards with arrows
steps = [
    ("1", "DXF Estrutural", "Engenheiro envia\nplanta de formas", ACCENT_BLUE, Inches(0.3)),
    ("2", "Escora.AI Processa", "Pipeline 6 estagios\n+ learning data", ACCENT_ORANGE, Inches(2.8)),
    ("3", "Output DXF", "Escoras posicionadas\n+ relatorio", ACCENT_GREEN, Inches(5.3)),
    ("4", "Engenheiro Revisa", "Ajusta posicoes,\nadiciona/remove", ACCENT_YELLOW, Inches(7.8)),
    ("5", "Feedback Loop", "Sistema compara\noriginal vs revisado", ACCENT_ORANGE, Inches(10.3)),
]

for num, title, desc, color, left in steps:
    top = Inches(1.8)
    shape = add_shape(slide, left, top, Inches(2.3), Inches(1.8), DARK_CARD, radius=0.05)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(2.3), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(14)
    p = tf.paragraphs[0]
    p.text = f"Passo {num}"
    p.font.size = Pt(11)
    p.font.color.rgb = color
    add_paragraph(tf, title, size=14, color=TEXT_WHITE, bold=True, space_before=4)
    add_paragraph(tf, desc, size=11, color=TEXT_LIGHT, space_before=6)

# Arrow between steps
for i in range(4):
    left = Inches(2.65 + i * 2.5)
    arrow = slide.shapes.add_textbox(left, Inches(2.5), Inches(0.3), Inches(0.3))
    set_text(arrow, "→", size=22, color=TEXT_MUTED, bold=True)

# What gets learned
card(slide, Inches(0.3), Inches(4.0), Inches(6.1), Inches(3.0),
     "O Que o Sistema Aprende com a Revisao",
     ["Escoras adicionadas pelo engenheiro = areas sub-escoradas",
      "Escoras removidas = areas sobre-escoradas ou redundantes",
      "Posicoes ajustadas = calibracao de esoacamento ideal",
      "Novas vigas/pilares marcados = falso-negativos na deteccao",
      "Elementos removidos = falso-positivos na classificacao",
      "Secoes corrigidas = ajuste de dimensoes detectadas",
      "Modelo de escora trocado = preferencia de catalogo"],
     title_color=ACCENT_GREEN, item_size=14)

card(slide, Inches(6.9), Inches(4.0), Inches(6.1), Inches(3.0),
     "Como o Aprendizado Melhora o Proximo Projeto",
     ["Layers conhecidos ganham prioridade na classificacao",
      "Secoes frequentes viram fallback quando texto nao detecta",
      "Esoacamento medio por tipo de viga calibra distribuicao",
      "Pe-direito historico evita valor padrao incorreto",
      "Falso-positivos recorrentes geram filtros automaticos",
      "Taxa de consumo (kg/m2) se aproxima do real",
      "A cada 10 projetos, precisao aumenta significativamente"],
     title_color=ACCENT_ORANGE, item_size=14)


# ═══════════════════════════════════════════
# SLIDE 7: Roadmap — Next Iterations
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.8))
set_text(tx, "ROADMAP — PROXIMAS ITERACOES", size=36, color=ACCENT_BLUE, bold=True)

# Phase 1: Short term
card(slide, Inches(0.3), Inches(1.3), Inches(4.0), Inches(5.6),
     "FASE 1 — Curto Prazo (1-2 semanas)",
     ["[FEEDBACK LOOP]",
      "  Diff engine: compara output vs revisao",
      "  Detecta escoras add/rem/mov pelo engenheiro",
      "  Salva ajustes como training data",
      "",
      "[DETECCAO]",
      "  Filtrar vigas de detalhe/corte (fora da planta)",
      "  Reduzir falsos-positivos de pilares (184 > ~30)",
      "  Resolver utilizacao >100% (escora mais forte)",
      "",
      "[OUTPUT]",
      "  DXF com blocks de escora (nao so circulos)",
      "  Tabela de quantidades por tipo de escora",
      "  Indicar contra-flecha no DXF"],
     title_color=ACCENT_GREEN, item_size=12)

# Phase 2: Medium term
card(slide, Inches(4.7), Inches(1.3), Inches(4.0), Inches(5.6),
     "FASE 2 — Medio Prazo (1-2 meses)",
     ["[CATALOGO REAL]",
      "  Integrar catalogo de equipamentos reais",
      "  Torres modulares (1000x1000 a 1540x1500mm)",
      "  Vigas metalicas VM130/VM80",
      "  Escoras ajustaveis Mecanor/ESC310/ESC450",
      "",
      "[MULTI-NIVEL]",
      "  Processar multiplos pavimentos por DXF",
      "  Segmentacao espacial por nivel (Tipo, Cob...)",
      "",
      "[ANTIPROJETO]",
      "  Taxa de consumo (kg/m3 > kg/m2)",
      "  Orcamento rapido por area amostral",
      "  Competitivo com ~10 kg/m2 referencia",
      "",
      "[DWG SUPPORT]",
      "  Conversao automatica DWG > DXF"],
     title_color=ACCENT_YELLOW, item_size=12)

# Phase 3: Long term
card(slide, Inches(9.1), Inches(1.3), Inches(4.0), Inches(5.6),
     "FASE 3 — Longo Prazo (3-6 meses)",
     ["[PROJETO EXECUTIVO]",
      "  Otimizacao por estoque disponivel",
      "  Distribuicao stock-aware (o que tem no patio)",
      "  BOM com codigos reais do cliente",
      "",
      "[METODOS CONSTRUTIVOS]",
      "  Alvenaria estrutural (sem pilares/vigas)",
      "  Steel deck (auto-portante ate 3-4.3m)",
      "  Laje nervurada, protendida, pre-moldada",
      "",
      "[AI AVANCADA]",
      "  Computer vision para ler DXF scanneado",
      "  LLM para interpretar notas do engenheiro",
      "  Modelo preditivo de taxa por tipologia",
      "",
      "[PLATAFORMA]",
      "  SaaS web (upload DXF, download resultado)",
      "  Dashboard com metricas por cliente",
      "  API para integracao com ERPs"],
     title_color=ACCENT_BLUE, item_size=12)


# ═══════════════════════════════════════════
# SLIDE 8: Acceleration Factors
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.8))
set_text(tx, "FATORES DE ACELERACAO DO APRENDIZADO", size=34, color=ACCENT_ORANGE, bold=True)

tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(0.5))
set_text(tx2, "O que pode contribuir para a melhoria acelerada do script", size=18, color=TEXT_LIGHT)

items = [
    ("Arquivos PE Revisados", "Projetos executivos JA revisados pelo engenheiro = ground truth perfeito. Cada arquivo com escoras posicionadas ensina o layout ideal.", ACCENT_GREEN),
    ("Volume de DXFs Diversos", "Quanto mais escritorios diferentes, mais layers e convencoes o sistema aprende. 50 arquivos > 10 arquivos.", ACCENT_GREEN),
    ("Feedback Estruturado", "Engenheiro marca: 'faltou escora aqui', 'esta a mais', 'viga nao detectada'. Structured diff vs free-form.", ACCENT_GREEN),
    ("Catalogo Real do Cliente", "Equipamentos reais com capacidades, pesos, precos. Permite calculo de taxa e BOM precisos.", ACCENT_YELLOW),
    ("Validacao Cruzada", "Comparar output do script com projetos executivos reais. Medir: escoras corretas/total, posicao media de erro.", ACCENT_YELLOW),
    ("Regras Explicitas do Engenheiro", "Regras tipo: 'nunca escora a menos de X do pilar', 'nesta regiao usar torre, nao escora'. Knowledge capture.", ACCENT_YELLOW),
    ("Test Suite com Ground Truth", "Para cada DXF, ter o output esperado (n_vigas, n_pilares, n_escoras). Regression testing automatico.", ACCENT_BLUE),
    ("Batch Processing", "Rodar 50 arquivos de uma vez, comparar resultados, identificar padroes de erro sistematicos.", ACCENT_BLUE),
    ("Metricas de Qualidade", "Score de deteccao (precision/recall), taxa de utilizacao media, desvio da referencia do engenheiro.", ACCENT_BLUE),
    ("Parceria com Locadora", "Acesso a projetos reais, feedback de campo, validacao por engenheiro senior = aceleracao 10x.", ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(items):
    col = i % 2
    row = i // 2
    left = Inches(0.3 + col * 6.4)
    top = Inches(1.6 + row * 1.12)
    shape = add_shape(slide, left, top, Inches(6.1), Inches(1.02), DARK_CARD, radius=0.03)
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Pt(8), top + Pt(10), Pt(22), Pt(22))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    set_text(badge, str(i + 1), size=11, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    # Text
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(38)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.font.bold = True
    add_paragraph(tf, desc, size=10, color=TEXT_MUTED, space_before=2)


# ═══════════════════════════════════════════
# SLIDE 9: Technical Architecture
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.8))
set_text(tx, "ARQUITETURA TECNICA", size=36, color=ACCENT_BLUE, bold=True)

# Left: Stack
card(slide, Inches(0.3), Inches(1.3), Inches(4.0), Inches(2.5),
     "Stack Tecnologico",
     ["Python 3.11+ — linguagem principal",
      "ezdxf — leitura/escrita DXF (todas as versoes)",
      "Shapely — geometria computacional (polygonize)",
      "Pydantic — modelos de dados tipados",
      "Typer + Rich — CLI com interface rica",
      "pytest — test suite (27 arquivos, 2200+ linhas)"],
     title_color=ACCENT_BLUE, item_size=13)

# Right: Key Algorithms
card(slide, Inches(4.7), Inches(1.3), Inches(4.0), Inches(2.5),
     "Algoritmos Chave",
     ["Beam detection: O(n2) segment pairing + merge",
      "Pillar detection: rect filter + cluster dedup",
      "Text classification: regex + proximity scoring",
      "Confidence: geometric * textual * agreement",
      "Slab derivation: Shapely polygonize + snap",
      "Shore distribution: grid + exclusion zones"],
     title_color=ACCENT_ORANGE, item_size=13)

# Right: Standards
card(slide, Inches(9.1), Inches(1.3), Inches(4.0), Inches(2.5),
     "Normas Aplicadas",
     ["NBR 15696:2009 — Formas e escoramentos",
      "NBR 6120:2019 — Acoes para calculo",
      "NBR 8800:2024 — Estruturas de aco",
      "gamma_concreto = 25.0 kN/m3",
      "gamma_f = 1.4 (majoracao)",
      "q_sobrecarga = 1.5 kN/m2"],
     title_color=ACCENT_GREEN, item_size=13)

# Bottom: Data flow
card(slide, Inches(0.3), Inches(4.2), Inches(12.5), Inches(2.8),
     "Fluxo de Dados",
     ["DXF File → stage_parse.py → ParseResult (textos, segmentos, retangulos, circulos, polylines)",
      "ParseResult → stage_segment.py → LevelSegment[] (agrupado por nivel/pavimento)",
      "LevelSegment → stage_classify.py → ClassifiedElement[] (vigas, pilares com scores 0-1)",
      "ClassifiedElement → stage_metadata.py → pe_direito, espessura, escala extraidos",
      "ClassifiedElement → stage_calculate.py → CalculationResult (escoras posicionadas + cargas)",
      "CalculationResult → stage_learn.py → LearningRecord salvo em data/learning.json",
      "",
      "Learning Store: 238 records | 11 beam layers | 9 pillar layers | secoes frequentes | pe-direito historico"],
     title_color=ACCENT_BLUE, item_size=12)


# ═══════════════════════════════════════════
# SLIDE 10: Closing
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), Inches(13.333), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_ORANGE
bar.line.fill.background()

tx = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
tf = set_text(tx, "ESCORA.AI", size=56, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_paragraph(tf, "Do DXF ao projeto de escoramento em segundos", size=24, color=ACCENT_ORANGE, space_before=12)

# Key stats
stats = [
    ("6 Estagios", "Pipeline completo"),
    ("39 Vigas", "Detectadas automaticamente"),
    ("412 Escoras", "Posicionadas com calculo"),
    ("238 Records", "De aprendizado acumulado"),
]
for i, (val, label) in enumerate(stats):
    left = Inches(1.5 + i * 2.8)
    metric_card(slide, left, Inches(3.5), Inches(2.3), Inches(1.3), val, label,
                [ACCENT_BLUE, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_YELLOW][i])

tx2 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1.0))
tf2 = set_text(tx2, "Proximo passo: feedback loop com engenheiro real", size=20, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
add_paragraph(tf2, "Cada revisao torna o sistema mais preciso — IA que aprende na pratica", size=16, color=TEXT_MUTED, space_before=8)

tx3 = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
set_text(tx3, "Raphael Lages | escora.ai | Marco 2026", size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ── Save ──
output_dir = "/Users/raphaellages/Desktop/escora-ai/output"
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "Escora.AI-Apresentacao.pptx")
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
